#!/usr/bin/env python3
"""
海南7天游攻略PPT生成器 - 详细版
包含具体餐厅、酒店、景点推荐
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

# 创建演示文稿
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# 配色方案
COLORS = {
    'primary': RGBColor(0, 150, 136),
    'secondary': RGBColor(255, 152, 0),
    'accent': RGBColor(255, 87, 34),
    'dark': RGBColor(33, 33, 33),
    'light': RGBColor(255, 255, 255),
    'hotel': RGBColor(63, 81, 181),
    'food': RGBColor(233, 30, 99),
    'spot': RGBColor(0, 150, 136),
}

def add_title_slide(prs, title, subtitle):
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = COLORS['primary']
    shape.line.fill.background()
    
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(12), Inches(1.5))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(60)
    p.font.bold = True
    p.font.color.rgb = COLORS['light']
    
    sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.2), Inches(12), Inches(1))
    tf = sub_box.text_frame
    p = tf.paragraphs[0]
    p.text = subtitle
    p.font.size = Pt(28)
    p.font.color.rgb = COLORS['light']
    
    return slide

def add_detail_slide(prs, day_num, title, spots, foods, hotels, color):
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    # 日期标签
    day_label = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.4), Inches(0.3), Inches(1), Inches(0.5))
    day_label.fill.solid()
    day_label.fill.fore_color.rgb = color
    day_label.line.fill.background()
    
    day_text = slide.shapes.add_textbox(Inches(0.4), Inches(0.33), Inches(1), Inches(0.45))
    tf = day_text.text_frame
    p = tf.paragraphs[0]
    p.text = f"DAY {day_num}"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = COLORS['light']
    p.alignment = PP_ALIGN.CENTER
    
    # 标题
    title_box = slide.shapes.add_textbox(Inches(1.6), Inches(0.25), Inches(11), Inches(0.6))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = COLORS['dark']
    
    # 三栏布局
    sections = [
        ("🏝️ 景点推荐", spots, COLORS['spot']),
        ("🍽️ 餐厅推荐", foods, COLORS['food']),
        ("🏨 住宿推荐", hotels, COLORS['hotel']),
    ]
    
    x_positions = [0.4, 4.5, 8.6]
    
    for i, (section_title, items, section_color) in enumerate(sections):
        x = x_positions[i]
        
        # 栏目标题
        header_bg = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(1), Inches(3.8), Inches(0.5))
        header_bg.fill.solid()
        header_bg.fill.fore_color.rgb = section_color
        header_bg.line.fill.background()
        
        header = slide.shapes.add_textbox(Inches(x), Inches(1.05), Inches(3.8), Inches(0.4))
        tf = header.text_frame
        p = tf.paragraphs[0]
        p.text = section_title
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = COLORS['light']
        p.alignment = PP_ALIGN.CENTER
        
        # 内容
        y_pos = 1.7
        for item in items:
            name_box = slide.shapes.add_textbox(Inches(x), Inches(y_pos), Inches(3.8), Inches(0.35))
            tf = name_box.text_frame
            p = tf.paragraphs[0]
            p.text = f"• {item['name']}"
            p.font.size = Pt(13)
            p.font.bold = True
            p.font.color.rgb = section_color
            
            if 'desc' in item:
                desc_box = slide.shapes.add_textbox(Inches(x + 0.15), Inches(y_pos + 0.35), Inches(3.65), Inches(0.3))
                tf = desc_box.text_frame
                p = tf.paragraphs[0]
                p.text = item['desc']
                p.font.size = Pt(10)
                p.font.color.rgb = COLORS['dark']
                y_pos += 0.35
            
            if 'price' in item:
                price_box = slide.shapes.add_textbox(Inches(x + 0.15), Inches(y_pos + 0.35), Inches(3.65), Inches(0.25))
                tf = price_box.text_frame
                p = tf.paragraphs[0]
                p.text = f"💰 {item['price']}"
                p.font.size = Pt(9)
                p.font.color.rgb = COLORS['secondary']
                y_pos += 0.3
            
            y_pos += 0.75
    
    return slide

print("🌴 正在生成海南7天游详细攻略PPT...")

# 封面
add_title_slide(prs, "海南7天深度游", "超详细攻略 · 餐厅·酒店·景点全推荐")
print("✅ 封面")

# Day 1 海口
add_detail_slide(prs, 1, "抵达海口 | 骑楼老街+夜市",
    [{'name': '骑楼老街', 'desc': '南洋风情建筑群，免费', 'price': '⭐⭐⭐⭐⭐'},
     {'name': '海口钟楼', 'desc': '地标建筑，拍照打卡', 'price': '免费'},
     {'name': '万绿园', 'desc': '城市公园，散步休闲', 'price': '免费'}],
    [{'name': '海大南门夜市', 'desc': '清补凉、海南粉、烧烤', 'price': '¥30-60'},
     {'name': '陈记汤饭', 'desc': '水巷口老店，辣汤饭', 'price': '¥25-40'},
     {'name': '椰语堂', 'desc': '清补凉连锁，料足', 'price': '¥20-35'}],
    [{'name': '海口万豪酒店', 'desc': '龙华区五星', 'price': '¥600-900'},
     {'name': '皇马假日酒店', 'desc': '骑楼老街旁', 'price': '¥200-350'},
     {'name': '7天连锁', 'desc': '多分店可选', 'price': '¥120-200'}],
    COLORS['primary'])
print("✅ Day 1")

# Day 2 海口
add_detail_slide(prs, 2, "海口文化游 | 博物馆+古迹",
    [{'name': '海南省博物馆', 'desc': '了解海南历史，免费', 'price': '⭐⭐⭐⭐'},
     {'name': '五公祠', 'desc': '古建筑群', 'price': '¥20'},
     {'name': '世纪大桥', 'desc': '夜景打卡', 'price': '免费'}],
    [{'name': '琼菜记忆', 'desc': '海南菜连锁', 'price': '¥60-100'},
     {'name': '原味主张', 'desc': '椰子鸡火锅', 'price': '¥80-120'},
     {'name': '老爸茶餐厅', 'desc': '本地早茶', 'price': '¥30-50'}],
    [{'name': '海口希尔顿', 'desc': '国贸商圈', 'price': '¥500-800'},
     {'name': '全季酒店', 'desc': '中山路店', 'price': '¥250-400'},
     {'name': '如家精选', 'desc': '明珠广场', 'price': '¥150-250'}],
    COLORS['primary'])
print("✅ Day 2")

# Day 3 文昌
add_detail_slide(prs, 3, "文昌→博鳌 | 航天+椰林",
    [{'name': '文昌航天城', 'desc': '中国航天发射场', 'price': '¥100'},
     {'name': '东郊椰林', 'desc': '万亩椰林海岸', 'price': '免费'},
     {'name': '博鳌论坛会址', 'desc': '亚洲论坛永久会址', 'price': '¥60'}],
    [{'name': '力哥地道海南菜', 'desc': '文昌鸡正宗', 'price': '¥70-100'},
     {'name': '椰子大观园餐厅', 'desc': '椰子宴特色', 'price': '¥60-90'},
     {'name': '博鳌海鲜大排档', 'desc': '新鲜实惠', 'price': '¥80-150'}],
    [{'name': '博鳌亚洲论坛酒店', 'desc': '五星海景', 'price': '¥800-1500'},
     {'name': '博鳌和悦酒店', 'desc': '精品酒店', 'price': '¥400-600'},
     {'name': '博鳌海景民宿', 'desc': '性价比高', 'price': '¥200-350'}],
    COLORS['secondary'])
print("✅ Day 3")

# Day 4 万宁
add_detail_slide(prs, 4, "万宁 | 冲浪圣地",
    [{'name': '日月湾', 'desc': '冲浪圣地', 'price': '免费'},
     {'name': '石梅湾', 'desc': '最美海湾', 'price': '免费'},
     {'name': '兴隆热带植物园', 'desc': '热带植物', 'price': '¥50'}],
    [{'name': '南洋风味', 'desc': '兴隆特色', 'price': '¥50-80'},
     {'name': '石梅湾海鲜', 'desc': '海边餐厅', 'price': '¥100-180'},
     {'name': '日月湾冲浪餐吧', 'desc': '西餐简餐', 'price': '¥60-100'}],
    [{'name': '石梅湾威斯汀', 'desc': '五星度假', 'price': '¥1000-1800'},
     {'name': '艾美度假酒店', 'desc': '石梅湾', 'price': '¥900-1600'},
     {'name': '日月湾青旅', 'desc': '冲浪主题', 'price': '¥100-200'}],
    COLORS['accent'])
print("✅ Day 4")

# Day 5 三亚
add_detail_slide(prs, 5, "万宁→三亚 | 热带风情",
    [{'name': '亚龙湾热带天堂', 'desc': '森林公园', 'price': '¥145'},
     {'name': '三亚湾椰梦长廊', 'desc': '日落打卡', 'price': '免费'},
     {'name': '第一市场', 'desc': '海鲜市场', 'price': '免费'}],
    [{'name': '春园海鲜广场', 'desc': '明码标价', 'price': '¥150-300'},
     {'name': '嗲嗲的椰子鸡', 'desc': '网红店', 'price': '¥100-150'},
     {'name': '第一市场小吃', 'desc': '抱罗粉水果', 'price': '¥30-60'}],
    [{'name': '三亚湾红树林', 'desc': '度假综合体', 'price': '¥400-700'},
     {'name': '大东海酒店', 'desc': '海景房', 'price': '¥300-500'},
     {'name': '三亚湾民宿', 'desc': '性价比高', 'price': '¥150-300'}],
    RGBColor(233, 30, 99))
print("✅ Day 5")

# Day 6 蜈支洲
add_detail_slide(prs, 6, "三亚海岛游 | 蜈支洲岛",
    [{'name': '蜈支洲岛', 'desc': '中国马尔代夫', 'price': '¥144+项目'},
     {'name': '情人桥', 'desc': '网红打卡', 'price': '含门票'},
     {'name': '观日岩', 'desc': '日出最佳点', 'price': '含门票'}],
    [{'name': '岛上自助餐厅', 'desc': '方便但贵', 'price': '¥100-150'},
     {'name': '后海渔村', 'desc': '下岛后吃', 'price': '¥60-100'},
     {'name': '三亚湾海鲜', 'desc': '晚上大餐', 'price': '¥150-300'}],
    [{'name': '亚龙湾瑞吉', 'desc': '奢华五星', 'price': '¥1500-3000'},
     {'name': '三亚湾喜来登', 'desc': '海景房', 'price': '¥500-900'},
     {'name': '大东海民宿', 'desc': '近海边', 'price': '¥150-300'}],
    RGBColor(233, 30, 99))
print("✅ Day 6")

# Day 7 返程
add_detail_slide(prs, 7, "三亚深度游 | 返程",
    [{'name': '南山文化旅游区', 'desc': '108米观音', 'price': '¥122'},
     {'name': '天涯海角', 'desc': '浪漫地标', 'price': '¥68'},
     {'name': '大小洞天', 'desc': '道教文化', 'price': '¥75'}],
    [{'name': '南山素斋', 'desc': '特色素食', 'price': '¥80-120'},
     {'name': '椰小鸡', 'desc': '王思聪同款', 'price': '¥80-120'},
     {'name': '机场简餐', 'desc': '返程前吃', 'price': '¥50-80'}],
    [{'name': '亚特兰蒂斯', 'desc': '海棠湾顶奢', 'price': '¥2000-5000'},
     {'name': '三亚湾酒店', 'desc': '近机场', 'price': '¥300-600'},
     {'name': '机场附近酒店', 'desc': '早班机首选', 'price': '¥150-300'}],
    RGBColor(233, 30, 99))
print("✅ Day 7")

# 保存
output_path = "/Users/shangguan/.qclaw/workspace/海南7天游攻略-详细版.pptx"
prs.save(output_path)

print(f"\n🎉 PPT生成完成！")
print(f"📁 保存位置：{output_path}")
print(f"📊 共 {len(prs.slides)} 页")
