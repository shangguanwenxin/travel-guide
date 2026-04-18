#!/usr/bin/env python3
"""
海南7天游攻略PPT生成器
图文结合，精美排版
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import requests
from io import BytesIO
from PIL import Image
import os

# 创建演示文稿
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# 配色方案
COLORS = {
    'primary': RGBColor(0, 150, 136),      # 青绿色 - 海洋
    'secondary': RGBColor(255, 152, 0),    # 橙色 - 阳光
    'accent': RGBColor(255, 87, 34),       # 深橙 - 活力
    'dark': RGBColor(33, 33, 33),          # 深灰 - 文字
    'light': RGBColor(255, 255, 255),      # 白色
    'bg_light': RGBColor(240, 248, 255),   # 浅蓝背景
}

# 图片URL（使用占位图服务，实际使用时可以替换为真实图片）
IMAGES = {
    'cover': 'https://images.unsplash.com/photo-1559827260-dc66d52bef19?w=1200',  # 海滩
    'haikou': 'https://images.unsplash.com/photo-1599571234909-29ed5d1321d6?w=800',  # 老街
    'wenchang': 'https://images.unsplash.com/photo-1558618666-fcd25c85cd64?w=800',  # 椰林
    'boao': 'https://images.unsplash.com/photo-1507525428034-b723cf961d3e?w=800',  # 海滩
    'wanning': 'https://images.unsplash.com/photo-1502680390469-be75c86b636f?w=800',  # 冲浪
    'sanya': 'https://images.unsplash.com/photo-1540202404-a2f29016b523?w=800',  # 三亚
    'island': 'https://images.unsplash.com/photo-1559827260-dc66d52bef19?w=800',  # 海岛
    'food': 'https://images.unsplash.com/photo-1555126634-323283e090fa?w=800',  # 美食
}

def add_gradient_background(slide, color1, color2):
    """添加渐变背景"""
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color1

def add_title_slide(prs, title, subtitle):
    """创建标题页"""
    slide_layout = prs.slide_layouts[6]  # 空白布局
    slide = prs.slides.add_slide(slide_layout)
    
    # 添加背景色块
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = COLORS['primary']
    shape.line.fill.background()
    
    # 添加装饰圆形
    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(9), Inches(-1), Inches(6), Inches(6))
    circle.fill.solid()
    circle.fill.fore_color.rgb = COLORS['secondary']
    circle.fill.fore_color.brightness = 0.3
    circle.line.fill.background()
    
    # 主标题
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(12), Inches(1.5))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(60)
    p.font.bold = True
    p.font.color.rgb = COLORS['light']
    p.alignment = PP_ALIGN.LEFT
    
    # 副标题
    sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.2), Inches(12), Inches(1))
    tf = sub_box.text_frame
    p = tf.paragraphs[0]
    p.text = subtitle
    p.font.size = Pt(28)
    p.font.color.rgb = COLORS['light']
    p.alignment = PP_ALIGN.LEFT
    
    return slide

def add_day_slide(prs, day_num, title, activities, image_url, highlight_color):
    """创建每日行程页"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    # 左侧色块背景
    left_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(0.3), prs.slide_height)
    left_bar.fill.solid()
    left_bar.fill.fore_color.rgb = highlight_color
    left_bar.line.fill.background()
    
    # 顶部装饰条
    top_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(0.15))
    top_bar.fill.solid()
    top_bar.fill.fore_color.rgb = highlight_color
    top_bar.line.fill.background()
    
    # 日期标签
    day_label = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.6), Inches(0.4), Inches(1.2), Inches(0.6))
    day_label.fill.solid()
    day_label.fill.fore_color.rgb = highlight_color
    day_label.line.fill.background()
    
    day_text = slide.shapes.add_textbox(Inches(0.6), Inches(0.45), Inches(1.2), Inches(0.5))
    tf = day_text.text_frame
    p = tf.paragraphs[0]
    p.text = f"DAY {day_num}"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = COLORS['light']
    p.alignment = PP_ALIGN.CENTER
    
    # 标题
    title_box = slide.shapes.add_textbox(Inches(2), Inches(0.4), Inches(8), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = COLORS['dark']
    
    # 图片区域（右侧）
    try:
        img_response = requests.get(image_url, timeout=10)
        if img_response.status_code == 200:
            img = Image.open(BytesIO(img_response.content))
            img_path = f'/tmp/hainan_img_{day_num}.jpg'
            img.save(img_path)
            
            # 添加图片框
            pic = slide.shapes.add_picture(img_path, Inches(7.5), Inches(1.5), width=Inches(5.5))
            # 添加图片边框装饰
            img_border = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(7.3), Inches(1.3), Inches(5.9), Inches(4.4))
            img_border.fill.background()
            img_border.line.color.rgb = highlight_color
            img_border.line.width = Pt(3)
    except:
        # 图片加载失败时添加占位符
        placeholder = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(7.5), Inches(1.5), Inches(5.5), Inches(4))
        placeholder.fill.solid()
        placeholder.fill.fore_color.rgb = COLORS['bg_light']
        placeholder.line.color.rgb = highlight_color
        placeholder.line.width = Pt(2)
    
    # 活动内容（左侧）
    content_box = slide.shapes.add_textbox(Inches(0.6), Inches(1.5), Inches(6.5), Inches(5.5))
    tf = content_box.text_frame
    tf.word_wrap = True
    
    for i, activity in enumerate(activities):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        
        # 时间点加粗
        if '：' in activity:
            parts = activity.split('：', 1)
            p.text = f"▸ {activity}"
        else:
            p.text = f"▸ {activity}"
        
        p.font.size = Pt(18)
        p.font.color.rgb = COLORS['dark']
        p.space_before = Pt(12)
        p.line_spacing = 1.3
    
    return slide

def add_overview_slide(prs):
    """创建行程概览页"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    # 标题
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "🗺️ 7天行程概览"
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = COLORS['primary']
    
    # 行程路线
    route = [
        ("Day 1-2", "海口", "骑楼老街 | 省博物馆 | 夜市", COLORS['primary']),
        ("Day 3", "文昌→博鳌", "航天城 | 东郊椰林 | 博鳌论坛", COLORS['secondary']),
        ("Day 4", "万宁", "日月湾冲浪 | 石梅湾", COLORS['accent']),
        ("Day 5-7", "三亚", "蜈支洲 | 南山 | 天涯海角", RGBColor(233, 30, 99)),
    ]
    
    y_pos = 1.5
    for day, city, spots, color in route:
        # 色块背景
        bg = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(y_pos), Inches(12), Inches(1.2))
        bg.fill.solid()
        bg.fill.fore_color.rgb = color
        bg.fill.fore_color.brightness = 0.85
        bg.line.fill.background()
        
        # 日期标签
        day_tag = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.7), Inches(y_pos + 0.25), Inches(1.3), Inches(0.7))
        day_tag.fill.solid()
        day_tag.fill.fore_color.rgb = color
        day_tag.line.fill.background()
        
        day_text = slide.shapes.add_textbox(Inches(0.7), Inches(y_pos + 0.3), Inches(1.3), Inches(0.6))
        tf = day_text.text_frame
        p = tf.paragraphs[0]
        p.text = day
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = COLORS['light']
        p.alignment = PP_ALIGN.CENTER
        
        # 城市名
        city_text = slide.shapes.add_textbox(Inches(2.2), Inches(y_pos + 0.2), Inches(3), Inches(0.8))
        tf = city_text.text_frame
        p = tf.paragraphs[0]
        p.text = city
        p.font.size = Pt(28)
        p.font.bold = True
        p.font.color.rgb = color
        
        # 景点
        spots_text = slide.shapes.add_textbox(Inches(5), Inches(y_pos + 0.35), Inches(7), Inches(0.6))
        tf = spots_text.text_frame
        p = tf.paragraphs[0]
        p.text = spots
        p.font.size = Pt(16)
        p.font.color.rgb = COLORS['dark']
        
        y_pos += 1.5
    
    return slide

def add_food_slide(prs):
    """创建美食推荐页"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    # 标题
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "🍜 必吃美食清单"
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = COLORS['accent']
    
    foods = [
        ("文昌鸡", "海南四大名菜之首", "文昌、海口"),
        ("加积鸭", "皮薄肉嫩", "琼海"),
        ("东山羊", "无膻味，肉质鲜美", "万宁"),
        ("和乐蟹", "膏满肉肥", "万宁"),
        ("海南粉", "早餐首选", "全岛"),
        ("清补凉", "消暑甜品", "全岛"),
        ("椰子鸡", "清甜养生", "三亚"),
        ("海鲜大餐", "现买现做", "三亚第一市场"),
    ]
    
    # 左列
    y_pos = 1.3
    for i, (name, desc, location) in enumerate(foods[:4]):
        # 美食名称
        name_box = slide.shapes.add_textbox(Inches(0.5), Inches(y_pos), Inches(3), Inches(0.5))
        tf = name_box.text_frame
        p = tf.paragraphs[0]
        p.text = f"🍽️ {name}"
        p.font.size = Pt(20)
        p.font.bold = True
        p.font.color.rgb = COLORS['accent']
        
        # 描述
        desc_box = slide.shapes.add_textbox(Inches(0.5), Inches(y_pos + 0.5), Inches(5.5), Inches(0.4))
        tf = desc_box.text_frame
        p = tf.paragraphs[0]
        p.text = f"{desc}  |  📍 {location}"
        p.font.size = Pt(14)
        p.font.color.rgb = COLORS['dark']
        
        y_pos += 1.4
    
    # 右列
    y_pos = 1.3
    for i, (name, desc, location) in enumerate(foods[4:]):
        name_box = slide.shapes.add_textbox(Inches(6.8), Inches(y_pos), Inches(3), Inches(0.5))
        tf = name_box.text_frame
        p = tf.paragraphs[0]
        p.text = f"🍽️ {name}"
        p.font.size = Pt(20)
        p.font.bold = True
        p.font.color.rgb = COLORS['accent']
        
        desc_box = slide.shapes.add_textbox(Inches(6.8), Inches(y_pos + 0.5), Inches(5.5), Inches(0.4))
        tf = desc_box.text_frame
        p = tf.paragraphs[0]
        p.text = f"{desc}  |  📍 {location}"
        p.font.size = Pt(14)
        p.font.color.rgb = COLORS['dark']
        
        y_pos += 1.4
    
    return slide

def add_tips_slide(prs):
    """创建实用信息页"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    # 标题
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "💡 实用出行贴士"
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = COLORS['primary']
    
    # 内容分三栏
    sections = [
        ("🌤️ 最佳旅行时间", [
            "11月-次年3月：气候最舒适",
            "避寒首选，阳光充足",
            "避开春节、国庆高峰期",
        ], COLORS['secondary']),
        ("🚗 交通建议", [
            "飞机：海口/三亚机场",
            "高铁：环岛高铁2小时",
            "自驾：最推荐！风景绝美",
        ], COLORS['primary']),
        ("🎒 必备物品", [
            "防晒霜 SPF50+",
            "墨镜、遮阳帽",
            "泳衣、沙滩鞋",
            "肠胃药（备用）",
        ], COLORS['accent']),
    ]
    
    x_positions = [0.5, 4.5, 8.5]
    for i, (title, items, color) in enumerate(sections):
        x = x_positions[i]
        
        # 栏目标题
        title_box = slide.shapes.add_textbox(Inches(x), Inches(1.3), Inches(3.8), Inches(0.6))
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(20)
        p.font.bold = True
        p.font.color.rgb = color
        
        # 内容
        y_pos = 2.0
        for item in items:
            item_box = slide.shapes.add_textbox(Inches(x), Inches(y_pos), Inches(3.8), Inches(0.5))
            tf = item_box.text_frame
            p = tf.paragraphs[0]
            p.text = f"• {item}"
            p.font.size = Pt(14)
            p.font.color.rgb = COLORS['dark']
            y_pos += 0.6
    
    # 预算参考
    budget_box = slide.shapes.add_textbox(Inches(0.5), Inches(5), Inches(12), Inches(2))
    tf = budget_box.text_frame
    
    p = tf.paragraphs[0]
    p.text = "💰 预算参考（人均7天）"
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = COLORS['dark']
    
    budgets = [
        ("经济型", "3000-4000元", COLORS['primary']),
        ("舒适型", "5000-7000元", COLORS['secondary']),
        ("豪华型", "8000元以上", COLORS['accent']),
    ]
    
    p = tf.add_paragraph()
    p.text = ""
    p.space_before = Pt(10)
    
    for name, amount, color in budgets:
        p = tf.add_paragraph()
        p.text = f"{name}: {amount}"
        p.font.size = Pt(16)
        p.font.color.rgb = color
        p.space_before = Pt(6)
    
    return slide

def add_end_slide(prs):
    """创建结束页"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    # 背景
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = COLORS['primary']
    bg.line.fill.background()
    
    # 装饰
    circle1 = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(-2), Inches(-2), Inches(6), Inches(6))
    circle1.fill.solid()
    circle1.fill.fore_color.rgb = COLORS['secondary']
    circle1.fill.fore_color.brightness = 0.2
    circle1.line.fill.background()
    
    circle2 = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(10), Inches(4), Inches(5), Inches(5))
    circle2.fill.solid()
    circle2.fill.fore_color.rgb = COLORS['accent']
    circle2.fill.fore_color.brightness = 0.3
    circle2.line.fill.background()
    
    # 主标题
    title_box = slide.shapes.add_textbox(Inches(0), Inches(2.5), prs.slide_width, Inches(1.5))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "🌴 海南，等你来！"
    p.font.size = Pt(54)
    p.font.bold = True
    p.font.color.rgb = COLORS['light']
    p.alignment = PP_ALIGN.CENTER
    
    # 副标题
    sub_box = slide.shapes.add_textbox(Inches(0), Inches(4.2), prs.slide_width, Inches(1))
    tf = sub_box.text_frame
    p = tf.paragraphs[0]
    p.text = "阳光、沙滩、美食，开启你的热带之旅"
    p.font.size = Pt(24)
    p.font.color.rgb = COLORS['light']
    p.alignment = PP_ALIGN.CENTER
    
    return slide

# ==================== 生成PPT ====================

print("🌴 正在生成海南7天游攻略PPT...")

# 1. 封面
add_title_slide(prs, "海南7天深度游", "阳光 · 沙滩 · 美食 · 热带风情")
print("✅ 封面页")

# 2. 行程概览
add_overview_slide(prs)
print("✅ 行程概览页")

# 3-9. 每日行程
days = [
    (1, "抵达海口 | 骑楼老街漫步", [
        "上午/中午：抵达海口美兰机场",
        "下午：骑楼老街，感受南洋风情",
        "晚上：海大南门夜市，品尝海南粉、清补凉",
        "住宿：海口市区",
    ], IMAGES['haikou'], COLORS['primary']),
    (2, "海口文化游 | 博物馆与古迹", [
        "上午：海南省博物馆，了解海南历史",
        "下午：五公祠 + 海口钟楼",
        "晚上：万绿园散步，世纪大桥看夜景",
        "住宿：海口市区",
    ], IMAGES['haikou'], COLORS['primary']),
    (3, "文昌→博鳌 | 航天与椰林", [
        "上午：文昌航天发射中心参观",
        "中午：品尝正宗文昌鸡",
        "下午：东郊椰林，椰子大观园",
        "傍晚：抵达博鳌，海边看日落",
        "住宿：博鳌镇",
    ], IMAGES['wenchang'], COLORS['secondary']),
    (4, "万宁 | 冲浪圣地", [
        "上午：博鳌亚洲论坛永久会址",
        "中午：博鳌小镇午餐",
        "下午：前往万宁日月湾",
        "傍晚：冲浪体验或海边漫步",
        "住宿：万宁石梅湾",
    ], IMAGES['wanning'], COLORS['accent']),
    (5, "万宁→三亚 | 热带风情", [
        "上午：兴隆热带植物园",
        "中午：品尝兴隆咖啡+东南亚风味",
        "下午：前往三亚（约1.5小时）",
        "晚上：第一市场海鲜大餐",
        "住宿：三亚湾",
    ], IMAGES['sanya'], RGBColor(233, 30, 99)),
    (6, "三亚海岛游 | 蜈支洲岛", [
        "上午：蜈支洲岛一日游",
        "下午：潜水/水上项目体验",
        "傍晚：返回三亚市区",
        "晚上：三亚湾椰梦长廊看日落",
        "住宿：三亚湾",
    ], IMAGES['island'], RGBColor(233, 30, 99)),
    (7, "三亚深度游 | 返程", [
        "上午：南山文化旅游区",
        "中午：南山素斋",
        "下午：天涯海角或亚龙湾",
        "傍晚：根据航班时间前往机场",
        "结束愉快的海南之旅！",
    ], IMAGES['sanya'], RGBColor(233, 30, 99)),
]

for day_num, title, activities, img_url, color in days:
    add_day_slide(prs, day_num, title, activities, img_url, color)
    print(f"✅ Day {day_num} 行程页")

# 10. 美食推荐
add_food_slide(prs)
print("✅ 美食推荐页")

# 11. 实用贴士
add_tips_slide(prs)
print("✅ 实用贴士页")

# 12. 结束页
add_end_slide(prs)
print("✅ 结束页")

# 保存PPT
output_path = "/Users/shangguan/.qclaw/workspace/海南7天游攻略.pptx"
prs.save(output_path)

print(f"\n🎉 PPT生成完成！")
print(f"📁 保存位置：{output_path}")
print(f"📊 共 {len(prs.slides)} 页")
